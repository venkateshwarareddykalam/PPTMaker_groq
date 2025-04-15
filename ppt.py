import datetime
import os
from groq import Groq
from pptx import Presentation
from pptx.util import Pt

client = Groq(api_key="your_groq_api_key_free_in_groq_website")
topic = input("Enter the presentation topic: ")
student_name = input("Enter the name: ")
roll_number = input("Enter the roll number: ")

def get_all_slides_content(topic, num_slides):
    """
    Call the Groq API once to generate content for slides 2 to num_slides.
    The AI outputs slide content separated by a string of 10 dollar symbols ($$$$$$$$$$).
    Each slide has a title (first line) followed by bullet points.
    """
    if num_slides < 2:
        return None  # No content slides to generate

    prompt = (
        f"Generate content for slides 2 to {num_slides} of a {num_slides}-slide presentation on '{topic}'. "
        f"This means you should generate content for exactly {num_slides - 1} slides. "
        "For each slide, output the slide data in plain text, "
        "where the first line is the slide title and 8 subsequent lines are bullet points. "
        "After the data for each slide, output a continuous string of 10 dollar symbols (i.e., $$$$$$$$$$) "
        "as a separator. Do not output any additional text. For example:\n\n"
        "Slide Title 2\n"
        "Bullet point 1\n"
        "Bullet point 2\n"
        "$$$$$$$$$$\n"
        "Slide Title 3\n"
        "Bullet point 1\n"
        "$$$$$$$$$$\n"
        "...\n"
        "\n"
        "Ensure that your output ends with the 10 dollar symbols separator after the last slide."
    )

    try:
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.3-70b-versatile",
            stream=False,
        )
        ai_text = chat_completion.choices[0].message.content.strip()
        print("Raw AI response:")
        print(ai_text)  # Debug: display the raw response
        return ai_text
    except Exception as e:
        print("Error calling or parsing Groq API response:", e)
        return None

def parse_slide_content(raw_text):
    """
    Parse the raw AI-generated text into a list of slides.
    Each slide is a dictionary with 'title' and 'bullets' keys.
    """
    slides_raw = raw_text.split("$$$$$$$$$$")
    slides = []
    for segment in slides_raw:
        segment = segment.strip()
        if not segment:
            continue
        lines = [line.strip() for line in segment.splitlines() if line.strip()]
        if not lines:
            continue
        slide_title = lines[0]
        slide_bullets = lines[1:]  # Bullet points (if any)
        slides.append({"title": slide_title, "bullets": slide_bullets})
    return slides

def set_font_size(shape, size_pt=14):
    """Set font size for all paragraphs in a shape."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)

def add_title_slide(prs, title, subtitle):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_font_size(slide.shapes.title, 14)
    set_font_size(slide.placeholders[1], 14)

def add_content_slide(prs, title, content_lines):
    """Add a content slide with a title and bullet points."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]

    title_shape.text = title
    tf = content_shape.text_frame
    tf.clear()

    if content_lines:
        p = tf.paragraphs[0]
        p.text = content_lines[0]
        for line in content_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
    else:
        tf.text = ""

    set_font_size(title_shape, 14)
    set_font_size(content_shape, 14)

def generate_presentation(topic, student_name, roll_number, num_slides):
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_title = f"Presentation on {topic}"
    title_slide_subtitle = f"{student_name}   {roll_number}"
    add_title_slide(prs, title_slide_title, title_slide_subtitle)
    print("Slide 1 generated: Title Slide")

    # Generate content slides if num_slides >= 2
    if num_slides >= 2:
        raw_ai_text = get_all_slides_content(topic, num_slides)
        if raw_ai_text is None:
            print("Failed to retrieve AI content. Only the title slide will be generated.")
        else:
            slides_data = parse_slide_content(raw_ai_text)
            if not slides_data:
                print("No slide content could be parsed. Only the title slide will be generated.")
            else:
                for i, slide in enumerate(slides_data, start=2):
                    slide_title = slide.get("title", f"Slide {i}")
                    slide_bullets = slide.get("bullets", [])
                    add_content_slide(prs, slide_title, slide_bullets)
                    print(f"Slide {i} generated: {slide_title}")
    else:
        print("Only the title slide will be generated.")
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    presentation_filename = f"{topic.replace(' ', '_')}_Presentation.pptx"
    prs.save(presentation_filename+timestamp)
    print(f"Presentation saved as '{presentation_filename}'.")


while True:
    try:
        num_slides = int(input("Enter the total number of slides (must be an integer >= 1): "))
        if num_slides >= 1:
            break
        else:
            print("Please enter an integer greater than or equal to 1.")
    except ValueError:
        print("Invalid input. Please enter an integer.")
generate_presentation(topic, student_name, roll_number, num_slides)